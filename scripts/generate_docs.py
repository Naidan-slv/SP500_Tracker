#!/usr/bin/env python3
"""Generate Technical Report PDF and Presentation PPTX for SP500 Tracker submission."""

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "docs")

# Path to a TTF font that supports Unicode. We'll try DejaVu (common on most systems).
def _find_unicode_font():
    """Find a Unicode-capable TTF font on the system."""
    candidates = [
        "/System/Library/Fonts/Helvetica.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return None

UNICODE_FONT_PATH = _find_unicode_font()


def _sanitise(text: str) -> str:
    """Replace Unicode characters that latin-1 can't encode."""
    return (
        text
        .replace("\u2014", " - ")   # em dash
        .replace("\u2013", "-")     # en dash
        .replace("\u2018", "'")     # left single quote
        .replace("\u2019", "'")     # right single quote
        .replace("\u201c", '"')     # left double quote
        .replace("\u201d", '"')     # right double quote
        .replace("\u2026", "...")   # ellipsis
        .replace("\u2192", "->")   # right arrow
        .replace("\u2190", "<-")   # left arrow
        .replace("\u2022", "-")    # bullet
        .replace("\u25b6", ">")    # right triangle
        .replace("\u2502", "|")    # box drawing
        .replace("\u250c", "+")
        .replace("\u2510", "+")
        .replace("\u2514", "+")
        .replace("\u2518", "+")
        .replace("\u2500", "-")
    )


# ═══════════════════════════════════════════════════════════════════════════════
# TECHNICAL REPORT PDF
# ═══════════════════════════════════════════════════════════════════════════════

class TechnicalReportPDF(FPDF):
    """Custom PDF with header/footer for the technical report."""

    def normalize_text(self, text):
        """Override to sanitise Unicode before FPDF's latin-1 encoding."""
        return super().normalize_text(_sanitise(text))

    def header(self):
        if self.page_no() == 1:
            return  # no header on title page
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 8, "SP500 Tracker — Technical Report", align="L")
        self.cell(0, 8, "COMP3011", align="R", new_x="LMARGIN", new_y="NEXT")
        self.line(self.l_margin, self.get_y(), self.w - self.r_margin, self.get_y())
        self.ln(4)

    def footer(self):
        if self.page_no() == 1:
            return
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 10, f"Page {self.page_no() - 1}", align="C")

    def section_title(self, number, title):
        self.set_font("Helvetica", "B", 14)
        self.set_text_color(25, 25, 112)
        self.cell(0, 10, f"{number}. {title}", new_x="LMARGIN", new_y="NEXT")
        self.line(self.l_margin, self.get_y(), self.w - self.r_margin, self.get_y())
        self.ln(3)

    def subsection_title(self, title):
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(50, 50, 50)
        self.cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")
        self.ln(1)

    def body_text(self, text):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(30, 30, 30)
        self.multi_cell(0, 5.5, text)
        self.ln(2)

    def bullet(self, text):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(30, 30, 30)
        x = self.get_x()
        self.cell(6, 5.5, "•")
        self.multi_cell(0, 5.5, text)
        self.ln(1)

    def bold_bullet(self, bold_part, rest):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(30, 30, 30)
        self.cell(6, 5.5, "•")
        self.set_font("Helvetica", "B", 10)
        self.write(5.5, bold_part)
        self.set_font("Helvetica", "", 10)
        self.write(5.5, rest)
        self.ln(6)

    def add_table(self, headers, rows, col_widths=None):
        if col_widths is None:
            usable = self.w - self.l_margin - self.r_margin
            col_widths = [usable / len(headers)] * len(headers)

        # Header row
        self.set_font("Helvetica", "B", 9)
        self.set_fill_color(230, 235, 245)
        self.set_text_color(30, 30, 30)
        for i, h in enumerate(headers):
            self.cell(col_widths[i], 7, h, border=1, fill=True, align="C")
        self.ln()

        # Data rows
        self.set_font("Helvetica", "", 9)
        fill = False
        for row in rows:
            if fill:
                self.set_fill_color(245, 247, 250)
            else:
                self.set_fill_color(255, 255, 255)
            max_h = 7
            for i, cell in enumerate(row):
                self.cell(col_widths[i], max_h, str(cell), border=1, fill=True, align="C" if i > 0 else "L")
            self.ln()
            fill = not fill
        self.ln(3)


def build_technical_report():
    pdf = TechnicalReportPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_left_margin(20)
    pdf.set_right_margin(20)

    # ── Title Page ────────────────────────────────────────────────────────────
    pdf.add_page()
    pdf.ln(45)
    pdf.set_font("Helvetica", "B", 28)
    pdf.set_text_color(25, 25, 112)
    pdf.cell(0, 14, "SP500 Tracker", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)
    pdf.set_font("Helvetica", "", 16)
    pdf.set_text_color(80, 80, 80)
    pdf.cell(0, 10, "Technical Report", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(20)

    pdf.set_font("Helvetica", "", 12)
    pdf.set_text_color(50, 50, 50)
    lines = [
        "Module: COMP3011 — Web-Based Application Development",
        "Student: Naidan Salvador",
        "Date: March 2025",
        "",
        "GitHub: github.com/Naidan-slv/SP500_Tracker",
        "Live API: sp500-tracker.onrender.com",
    ]
    for line in lines:
        pdf.cell(0, 8, line, align="C", new_x="LMARGIN", new_y="NEXT")

    pdf.ln(30)
    pdf.set_font("Helvetica", "I", 10)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 8, "Word Count: ~2,400 (excluding tables, diagrams, and references)", align="C", new_x="LMARGIN", new_y="NEXT")

    # ── Page 1: Introduction + Stack Justification ────────────────────────────
    pdf.add_page()

    pdf.section_title("1", "Introduction")
    pdf.body_text(
        "The SP500 Tracker is a full-stack web application that enables users to browse S&P 500 stocks, "
        "view historical and live price charts, track watchlists with computed analytics, and manage investment "
        "portfolios. The application consumes a Kaggle dataset of historical OHLCV data for the top 50 S&P 500 "
        "companies (Shahrukh, 2025) and enriches it with live market data from Finnhub (Finnhub, 2025) and "
        "Yahoo Finance. This report details the technology choices, architectural design, development challenges, "
        "testing strategy, and areas for future improvement."
    )

    pdf.section_title("2", "Technology Stack Justification")

    pdf.subsection_title("2.1 Backend — FastAPI + SQLAlchemy")
    pdf.body_text(
        "FastAPI (Ramírez, 2018) was chosen over Flask or Django for its automatic OpenAPI documentation, "
        "Pydantic-based request/response validation, and native async/await support — critical for a service "
        "that aggregates data from multiple external APIs concurrently. SQLAlchemy 2.0 (Bayer, 2025) with the "
        "modern select() API provides type-safe query construction, while Alembic handles versioned schema "
        "migrations. PostgreSQL (via Supabase) was selected for production because it handles concurrent "
        "connections and scales to the 230,000+ OHLCV price rows stored in the application."
    )

    pdf.subsection_title("2.2 Frontend — React + TypeScript + Vite")
    pdf.body_text(
        "React 18 (Meta, 2025) was chosen for its component model and ecosystem maturity. TypeScript catches "
        "type errors at compile time, which proved valuable when modelling the 23-endpoint API surface with "
        "strongly-typed request/response interfaces. Vite (You, 2025) provides sub-second hot module replacement "
        "and optimised production builds. TanStack Query v5 (Linsley, 2025) manages all server state — automatic "
        "caching, background refetching, and stale-while-revalidate — eliminating manual useEffect data-fetching "
        "patterns. Recharts (Recharts Contributors, 2025) was selected for its composable React-native chart "
        "components."
    )

    pdf.subsection_title("2.3 Authentication — JWT")
    pdf.body_text(
        "JSON Web Tokens via python-jose (Meza, 2023) provide stateless authentication. Tokens are stored in "
        "the frontend's localStorage and attached to requests via the Authorization: Bearer header. This "
        "approach avoids server-side session storage and simplifies horizontal scaling (Jones, Bradley and "
        "Sakimura, 2015)."
    )

    pdf.subsection_title("2.4 External Data Providers")
    pdf.body_text(
        "A multi-provider cascade ensures data resilience. Finnhub (primary) provides real-time candles and "
        "company profiles. Yahoo Finance serves as a secondary fallback. Stored database OHLCV data acts as "
        "a guaranteed tertiary fallback. Google News RSS provides ticker-specific news with time-scoped queries. "
        "This cascade ensures the application always renders meaningful data even when individual providers fail."
    )

    # ── Page 2: Design & Architecture ─────────────────────────────────────────
    pdf.section_title("3", "Design & Architecture")

    pdf.subsection_title("3.1 System Architecture")
    pdf.body_text(
        "The application follows a three-tier architecture: a React single-page application communicates over "
        "HTTPS with a FastAPI REST API, which queries a PostgreSQL database and three external data providers "
        "(Finnhub, Yahoo Finance, Google News). Both the frontend (static site) and backend (web service) are "
        "deployed on Render (Render, 2025), with the database hosted on Supabase."
    )

    pdf.set_font("Courier", "", 8)
    pdf.set_text_color(60, 60, 60)
    arch_diagram = (
        "  +---------------+     HTTPS/JSON     +---------------+     SQL      +-------------+\n"
        "  |  React SPA    | ==================> |  FastAPI API  | ==========> | PostgreSQL  |\n"
        "  |  (Render       | <================  |  (Render Web  | <========== | (Supabase)  |\n"
        "  |   Static Site) |                    |   Service)    |             |             |\n"
        "  +---------------+                    +-------+-------+             +-------------+\n"
        "                                               |\n"
        "                                    +----------+----------+\n"
        "                                    |          |          |\n"
        "                                 Finnhub    Yahoo     Google\n"
        "                                 (Live)   (Fallback)  (News)\n"
    )
    pdf.multi_cell(0, 4, arch_diagram)
    pdf.ln(2)

    pdf.subsection_title("3.2 Backend Structure")
    pdf.body_text(
        "The backend uses a modular router pattern with four route modules: auth.py (3 endpoints for "
        "registration, login, and user profile), stocks.py (5 endpoints for search, detail, history, news, "
        "and live data), watchlists.py (7 endpoints for CRUD operations plus analytics/insights), and "
        "portfolios.py (8 endpoints for CRUD operations plus holdings management). Database models are "
        "defined in models.py using SQLAlchemy ORM: User, Stock, StockPrice, Watchlist, WatchlistItem, "
        "Portfolio, and PortfolioHolding."
    )

    pdf.subsection_title("3.3 Frontend Structure")
    pdf.body_text(
        "The frontend follows a page-component-hook pattern. Pages (src/pages/) orchestrate data fetching "
        "and layout. Components (src/components/) provide reusable UI elements such as charts, modals, and "
        "navigation. A centralised API layer (src/lib/api.ts) encapsulates all HTTP calls with typed "
        "request/response functions. AuthContext provides login state and token management via React context."
    )

    pdf.subsection_title("3.4 Caching Strategy")
    pdf.body_text(
        "A two-tier caching approach minimises external API calls and improves perceived performance:"
    )
    pdf.add_table(
        ["Layer", "Implementation", "TTL"],
        [
            ["Backend (server)", "Python dict (_LIVE_CACHE, _NEWS_CACHE)", "45s live, 5min news"],
            ["Frontend (client)", "TanStack Query v5", "30s–10min per resource"],
        ],
        col_widths=[40, 75, 55],
    )

    # ── Page 3: Data Source, Challenges & Testing ─────────────────────────────
    pdf.section_title("4", "Data Source")
    pdf.body_text(
        "The application's historical stock data is sourced from the 'Top 50 S&P 500 Companies' Kaggle "
        "dataset (Shahrukh, 2025), which provides daily OHLCV (Open, High, Low, Close, Volume) data for "
        "49 S&P 500 tickers spanning 2006–2026. This dataset was loaded into a PostgreSQL database comprising "
        "230,111 price rows across the stocks and stock_prices tables. Live and intraday data supplements "
        "this historical foundation via the Finnhub and Yahoo Finance APIs."
    )

    pdf.section_title("5", "Challenges & Lessons Learned")

    pdf.subsection_title("5.1 External API Reliability")
    pdf.body_text(
        "The most significant challenge was external API instability. Finnhub's free-tier rate limit "
        "(60 calls/min) and Yahoo Finance's unofficial API occasionally returning empty responses caused "
        "the live chart to break in production. The solution was a three-layer fallback cascade "
        "(Finnhub → Yahoo → Database history) with stale cache as a final safety net, ensuring the chart "
        "always renders meaningful data."
    )

    pdf.subsection_title("5.2 Fuzzy Search Resolution")
    pdf.body_text(
        'Supporting both ticker symbols and company names as input (e.g., typing "Berkshire" should '
        "resolve to BRK.B) required a normalisation pipeline that strips punctuation, spaces, and special "
        "characters before performing SQL LIKE queries. A hardcoded COMPANY_NAME_OVERRIDES map handles "
        "edge cases where database names do not match common queries."
    )

    pdf.subsection_title("5.3 CSS Layering & Overflow")
    pdf.body_text(
        "The autocomplete suggestion dropdown had z-index stacking conflicts with adjacent cards. Resolving "
        "this required managing overflow, position, and z-index across the layout shell, grid containers, "
        "and individual cards — a reminder that CSS stacking contexts must be considered holistically."
    )

    pdf.section_title("6", "Testing Approach")

    pdf.subsection_title("6.1 Strategy")
    pdf.body_text(
        "The project has 127 automated tests using pytest with an in-memory SQLite database for "
        "complete isolation. Tests are organised into focused modules:"
    )
    pdf.add_table(
        ["Module", "Tests", "Coverage"],
        [
            ["Auth (register, login, /me)", "15", "Registration, login, token validation, error cases"],
            ["Stocks (list, detail, history, news, live)", "20", "Search, pagination, timeframe validation"],
            ["Watchlists (CRUD + insights)", "11", "Create, list, delete, add/remove, duplicates"],
            ["Portfolios (CRUD + holdings)", "42", "Full CRUD, holdings edge cases"],
            ["Contract (API schema validation)", "10", "Response structure matches Pydantic models"],
            ["E2E (user journey)", "7", "Register → login → watchlist → insights"],
            ["Insights analytics", "21", "Volatility, change %, summary labels"],
        ],
        col_widths=[55, 15, 100],
    )

    pdf.subsection_title("6.2 Key Testing Decisions")
    pdf.bold_bullet("In-memory SQLite — ", "each test function gets a fresh database via pytest fixtures, ensuring complete isolation without network dependencies.")
    pdf.bold_bullet("Dependency injection override — ", "FastAPI's app.dependency_overrides replaces the production database session with a test session, avoiding any production data access.")
    pdf.bold_bullet("No external API mocking — ", "stock data endpoints operate on seeded test data; external providers (Finnhub, Yahoo) are not called during tests.")

    # ── Page 4: Limitations, Future Work, GenAI ───────────────────────────────
    pdf.section_title("7", "Limitations & Future Improvements")

    pdf.subsection_title("7.1 Current Limitations")
    pdf.bold_bullet("No real-time WebSocket streaming — ", "the live chart polls via REST; true real-time push would reduce latency.")
    pdf.bold_bullet("In-memory caching only — ", "the Python dict cache is not shared across multiple server instances; Redis would support horizontal scaling.")
    pdf.bold_bullet("No refresh token rotation — ", "JWT tokens expire after 60 minutes with no silent refresh mechanism.")
    pdf.bold_bullet("Equal-weight insights only — ", "watchlist analytics assume equal weighting; market-cap or custom weighting would be more analytically useful.")

    pdf.subsection_title("7.2 Future Improvements")
    pdf.bullet("WebSocket integration for push-based live price updates.")
    pdf.bullet("Redis caching to support multi-instance deployments.")
    pdf.bullet("OAuth 2.0 social login (Google, GitHub) alongside email/password.")
    pdf.bullet("Portfolio performance tracking — historical value charting, P&L calculations, benchmark comparison.")
    pdf.bullet("Mobile-responsive redesign using CSS container queries.")
    pdf.bullet("Alerting system — price threshold notifications via email or push.")

    pdf.section_title("8", "GenAI Declaration & Analysis")

    pdf.subsection_title("8.1 Tools Used")
    pdf.add_table(
        ["Tool", "Purpose"],
        [
            ["GitHub Copilot (Claude)", "Code generation, debugging, architecture, testing, documentation"],
        ],
        col_widths=[55, 115],
    )

    pdf.subsection_title("8.2 How GenAI Was Used")
    pdf.body_text(
        "Generative AI was used as a pair programming assistant throughout the project:"
    )
    pdf.bold_bullet("Code scaffolding — ", "initial endpoint structures, Pydantic models, and SQLAlchemy schemas were generated and then reviewed/modified by the developer.")
    pdf.bold_bullet("Test generation — ", "test cases were collaboratively designed, with AI generating boilerplate and the developer specifying edge cases and assertions.")
    pdf.bold_bullet("Debugging — ", "CSS stacking issues, external API integration failures, and database query optimisation were iteratively debugged with AI assistance.")
    pdf.bold_bullet("Documentation — ", "API documentation, this technical report, and the development log were drafted with AI and refined by the developer.")

    pdf.subsection_title("8.3 Critical Analysis")
    pdf.body_text(
        "Benefits: GenAI significantly accelerated development velocity — the 23-endpoint API with 127 tests "
        "was built faster than would have been possible manually. It caught subtle bugs (e.g., timezone-naive "
        "datetime comparisons, potential SQL injection vectors in LIKE queries) that might have been missed in "
        "manual review, and provided consistent code style across the codebase.\n\n"
        "Limitations: AI-generated CSS fixes sometimes addressed symptoms rather than root causes, requiring "
        "multiple iterations. Generated test cases occasionally tested implementation details rather than "
        "behaviour. All suggestions required careful review — none were blindly accepted.\n\n"
        "Conclusion: GenAI was a valuable productivity tool but required continuous human oversight. The "
        "developer maintained architectural decision-making authority while leveraging AI for implementation "
        "speed. All generated code was reviewed, tested, and understood before integration. A full interaction "
        "log is maintained in AI_DEVELOPMENT_LOG.md in the project repository."
    )

    # ── Page 5: References ────────────────────────────────────────────────────
    pdf.section_title("9", "References")
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(30, 30, 30)

    references = [
        "Bayer, M. (2025) SQLAlchemy 2.0 Documentation. Available at: https://docs.sqlalchemy.org/en/20/ (Accessed: 9 March 2025).",
        "Finnhub (2025) Finnhub Stock API Documentation. Available at: https://finnhub.io/docs/api (Accessed: 9 March 2025).",
        "Jones, M., Bradley, J. and Sakimura, N. (2015) 'JSON Web Token (JWT)', RFC 7519. Available at: https://datatracker.ietf.org/doc/html/rfc7519 (Accessed: 9 March 2025).",
        "Linsley, T. (2025) TanStack Query v5 Documentation. Available at: https://tanstack.com/query/latest (Accessed: 9 March 2025).",
        "Meta (2025) React Documentation. Available at: https://react.dev/ (Accessed: 9 March 2025).",
        "Meza, M. (2023) python-jose: A JOSE implementation in Python. Available at: https://github.com/mpdavis/python-jose (Accessed: 9 March 2025).",
        "Ramírez, S. (2018) FastAPI Documentation. Available at: https://fastapi.tiangolo.com/ (Accessed: 9 March 2025).",
        "Recharts Contributors (2025) Recharts: A composable charting library. Available at: https://recharts.org/ (Accessed: 9 March 2025).",
        "Render (2025) Render Cloud Application Hosting. Available at: https://render.com/ (Accessed: 9 March 2025).",
        "Shahrukh, I. (2025) Top 50 S&P 500 Companies Dataset. Available at: https://www.kaggle.com/datasets/ibrahimshahrukh/top-50-companies-dataset (Accessed: 9 March 2025).",
        "You, E. (2025) Vite: Next Generation Frontend Tooling. Available at: https://vitejs.dev/ (Accessed: 9 March 2025).",
    ]

    for ref in references:
        pdf.multi_cell(0, 5, ref)
        pdf.ln(2)

    # ── Save ──────────────────────────────────────────────────────────────────
    output_path = os.path.join(OUTPUT_DIR, "SP500_Tracker_Technical_Report.pdf")
    pdf.output(output_path)
    print(f"✅ Technical Report PDF saved to: {output_path}")
    return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# PRESENTATION PPTX
# ═══════════════════════════════════════════════════════════════════════════════

# Colour palette
BG_DARK    = RGBColor(18, 18, 40)
ACCENT     = RGBColor(25, 25, 112)
WHITE      = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 210)
MUTED      = RGBColor(160, 160, 175)
TABLE_HDR  = RGBColor(40, 40, 90)
TABLE_ALT  = RGBColor(30, 30, 60)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, left, top, width, height, font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, font_size=16, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    return tf


def add_bullets(slide, items, left, top, width, height, font_size=14, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf


def add_subtitle_line(slide, text, top_offset):
    """Add a coloured accent line under a subtitle."""
    add_title_text(slide, text, Inches(0.8), top_offset, Inches(8.4), Inches(0.5),
                   font_size=22, bold=True, color=RGBColor(100, 140, 255))


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "SP500 Tracker", Inches(1), Inches(1.8), Inches(11), Inches(1),
                   font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_title_text(slide, "A Full-Stack S&P 500 Stock Tracking Application", Inches(1), Inches(2.9), Inches(11), Inches(0.6),
                   font_size=20, bold=False, color=MUTED, alignment=PP_ALIGN.CENTER)
    add_body_text(slide,
                  "COMP3011 — Web-Based Application Development\n"
                  "Naidan Salvador  •  March 2025\n\n"
                  "GitHub: github.com/Naidan-slv/SP500_Tracker\n"
                  "Live API: sp500-tracker.onrender.com",
                  Inches(1), Inches(4.0), Inches(11), Inches(2.5),
                  font_size=16, color=LIGHT_GRAY)

    # ── Slide 2: Project Overview ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Project Overview", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)
    add_bullets(slide, [
        "🔍  Browse & search all S&P 500 stocks with fuzzy search",
        "📈  View historical OHLCV charts with configurable timeframes (1W – 5Y)",
        "⚡  Monitor live prices with intraday candlestick data",
        "📰  Read latest news aggregated from Google News RSS",
        "📋  Create watchlists with computed analytics & insights",
        "💼  Manage portfolios with holdings, cost tracking & allocation pie chart",
        "🔐  User authentication via JWT (register → login)",
        "",
        "23 API endpoints  •  127 automated tests  •  Deployed on Render",
    ], Inches(0.8), Inches(1.5), Inches(11), Inches(5), font_size=18, color=LIGHT_GRAY)

    # ── Slide 3: Technology Stack ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Technology Stack", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    stack_items = [
        ("Backend", "FastAPI + Python 3.12", "Auto-docs, async, Pydantic validation"),
        ("Database", "PostgreSQL (Supabase)", "230k+ price rows, concurrent access"),
        ("ORM", "SQLAlchemy 2.0 + Alembic", "Type-safe queries, versioned migrations"),
        ("Frontend", "React 18 + TypeScript + Vite", "Component model, type safety, fast HMR"),
        ("State Mgmt", "TanStack Query v5", "Caching, background refetch, stale-while-revalidate"),
        ("Charts", "Recharts", "React-native composable chart components"),
        ("Auth", "JWT (python-jose)", "Stateless, scalable, no session storage"),
        ("Deployment", "Render", "Backend web service + frontend static site"),
    ]

    # Build a table
    rows_count = len(stack_items) + 1
    cols_count = 3
    table_shape = slide.shapes.add_table(rows_count, cols_count,
                                          Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
    table = table_shape.table

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(5.0)

    # Header
    for j, hdr in enumerate(["Layer", "Technology", "Justification"]):
        cell = table.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER

    for i, (layer, tech, why) in enumerate(stack_items):
        for j, val in enumerate([layer, tech, why]):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ALT if i % 2 == 0 else BG_DARK
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.alignment = PP_ALIGN.LEFT if j > 0 else PP_ALIGN.CENTER

    # ── Slide 4: Architecture ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "System Architecture", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    arch_text = (
        "Three-tier architecture with multi-provider fallback:\n\n"
        "React SPA  ──HTTPS/JSON──▶  FastAPI API  ──SQL──▶  PostgreSQL (Supabase)\n"
        "  (Render Static)                (Render Web Service)\n\n"
        "                         External Data Providers:\n"
        "                     Finnhub (primary)  →  Yahoo Finance (fallback)  →  DB History (guaranteed)\n"
        "                                  Google News RSS (news articles)"
    )
    add_body_text(slide, arch_text, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3), font_size=16, color=LIGHT_GRAY)

    add_bullets(slide, [
        "Two-tier caching: backend in-memory (45s live, 5min news) + frontend TanStack Query (30s–10min)",
        "Stateless JWT authentication — no server-side sessions",
        "Infrastructure-as-Code deployment via render.yaml",
    ], Inches(0.8), Inches(4.5), Inches(11), Inches(2.5), font_size=16, color=MUTED)

    # ── Slide 5: API Overview ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "API Overview — 23 RESTful Endpoints", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    api_groups = [
        ("Auth (3)", "Register, Login, Me", "JWT tokens, password hashing (bcrypt)"),
        ("Stocks (5)", "List, Detail, History, News, Live", "Fuzzy search, OHLCV, 3-provider cascade"),
        ("Watchlists (7)", "CRUD + Items + Insights", "Volatility, % change, top gainer/loser"),
        ("Portfolios (8)", "CRUD + Holdings CRUD", "Quantity, avg cost, allocation breakdown"),
    ]

    rows_count = len(api_groups) + 1
    table_shape = slide.shapes.add_table(rows_count, 3,
                                          Inches(0.8), Inches(1.5), Inches(11.5), Inches(3))
    table = table_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(4.5)

    for j, hdr in enumerate(["Group", "Endpoints", "Key Features"]):
        cell = table.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE

    for i, (group, endpoints, features) in enumerate(api_groups):
        for j, val in enumerate([group, endpoints, features]):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ALT if i % 2 == 0 else BG_DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = LIGHT_GRAY

    add_bullets(slide, [
        "Interactive Swagger UI at /docs",
        "Full API documentation with example curl requests and JSON responses",
        "Consistent error handling with standard HTTP status codes",
    ], Inches(0.8), Inches(4.8), Inches(11), Inches(2), font_size=15, color=MUTED)

    # ── Slide 6: Key Features ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Key Features", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    add_bullets(slide, [
        '🔍  Fuzzy Stock Search — type "Berkshire" → resolves to BRK.B',
        "       SQL normalisation + override map + Yahoo Finance search fallback",
        "",
        "⚡  Live Chart with 3-Layer Fallback",
        "       Finnhub → Yahoo Finance → Database History → Stale Cache",
        "",
        "📋  Watchlist Insights — per-ticker analytics dashboard",
        "       1W/1M/1Y price change, 30-day volatility, average volume, top gainer/loser",
        "",
        "💼  Portfolio Allocation Pie Chart",
        "       Interactive donut chart with hover highlighting, total value in centre",
        "",
        "📰  Time-scoped News — Google News RSS with when:Xd operator",
    ], Inches(0.8), Inches(1.4), Inches(11), Inches(5.5), font_size=16, color=LIGHT_GRAY)

    # ── Slide 7: Version Control & Deployment ────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Version Control & Deployment", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    add_subtitle_line(slide, "Git Practices", Inches(1.3))
    add_bullets(slide, [
        "Conventional commit messages (feat:, fix:, docs:)",
        "Frequent, atomic commits — feature-by-feature",
        "Full history from project scaffold to final fixes",
    ], Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), font_size=15, color=LIGHT_GRAY)

    add_subtitle_line(slide, "Deployment", Inches(4.2))
    add_bullets(slide, [
        "render.yaml Infrastructure-as-Code configuration",
        "Backend: Python web service (Gunicorn + Uvicorn workers)",
        "Frontend: Static site with SPA redirect rules",
        "Environment variables managed via Render dashboard",
        "",
        "Live API: sp500-tracker.onrender.com",
        "Swagger Docs: sp500-tracker.onrender.com/docs",
    ], Inches(0.8), Inches(4.7), Inches(11), Inches(2.5), font_size=15, color=LIGHT_GRAY)

    # ── Slide 8: Testing ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Testing Strategy — 127 Automated Tests", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    test_data = [
        ("Auth", "15", "Register, login, token validation"),
        ("Stocks", "20", "Search, detail, history, pagination"),
        ("Watchlists", "11", "CRUD, duplicate prevention"),
        ("Portfolios", "42", "Full CRUD + holdings edge cases"),
        ("Contract", "10", "Response schema validation"),
        ("E2E Journey", "7", "Full user flow: register → insights"),
        ("Insights", "21", "Analytics accuracy, summary labels"),
    ]

    rows_count = len(test_data) + 1
    table_shape = slide.shapes.add_table(rows_count, 3,
                                          Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.2))
    table = table_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(7.5)

    for j, hdr in enumerate(["Category", "Count", "Description"]):
        cell = table.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE

    for i, (cat, count, desc) in enumerate(test_data):
        for j, val in enumerate([cat, count, desc]):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ALT if i % 2 == 0 else BG_DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = LIGHT_GRAY
                p.alignment = PP_ALIGN.CENTER if j == 1 else PP_ALIGN.LEFT

    add_bullets(slide, [
        "In-memory SQLite for complete isolation  •  FastAPI dependency injection overrides",
        "No external API calls during tests  •  All 127 tests pass in ~4 seconds",
    ], Inches(0.8), Inches(5.9), Inches(11), Inches(1.2), font_size=14, color=MUTED)

    # ── Slide 9: Challenges & GenAI ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Challenges & GenAI Usage", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    add_subtitle_line(slide, "Key Challenges", Inches(1.3))
    add_bullets(slide, [
        "1. External API instability → solved with 3-layer fallback + stale cache",
        "2. Fuzzy search edge cases (BRK.B, special chars) → normalisation pipeline",
        "3. CSS stacking conflicts → z-index management across layout layers",
    ], Inches(0.8), Inches(1.8), Inches(11), Inches(2), font_size=16, color=LIGHT_GRAY)

    add_subtitle_line(slide, "GenAI Usage (GitHub Copilot — Claude)", Inches(3.8))
    add_bullets(slide, [
        "Used as a pair programming assistant throughout the project",
        "Code scaffolding, test generation, debugging, documentation",
        "All code reviewed and understood before integration",
        "AI accelerated velocity but required continuous human oversight",
        "AI occasionally addressed symptoms over root causes (CSS issues needed 2 rounds)",
        "Full interaction log maintained in AI_DEVELOPMENT_LOG.md (19 entries)",
    ], Inches(0.8), Inches(4.3), Inches(11), Inches(3), font_size=15, color=LIGHT_GRAY)

    # ── Slide 10: Limitations & Future Work ──────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Limitations & Future Work", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    add_subtitle_line(slide, "Current Limitations", Inches(1.3))
    add_bullets(slide, [
        "REST polling for live data (no WebSocket push)",
        "In-memory cache not shared across server instances",
        "No refresh token rotation (60-min expiry)",
        "Equal-weight watchlist analytics only",
    ], Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), font_size=15, color=LIGHT_GRAY)

    add_subtitle_line(slide, "Future Improvements", Inches(4.2))
    add_bullets(slide, [
        "⚡  WebSocket real-time price streaming",
        "🗄️  Redis caching for horizontal scaling",
        "🔐  OAuth 2.0 social login (Google, GitHub)",
        "📊  Portfolio P&L tracking with benchmark comparison",
        "📱  Mobile-responsive redesign",
        "🔔  Price alerts via notifications",
    ], Inches(0.8), Inches(4.7), Inches(11), Inches(2.5), font_size=15, color=LIGHT_GRAY)

    # ── Slide 11: Summary ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_text(slide, "Summary & Deliverables", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), font_size=32)

    deliverables = [
        ("Working API (23 endpoints)", "Deployed on Render"),
        ("React Frontend (4 pages)", "Deployed on Render"),
        ("API Documentation", "Complete with examples"),
        ("Technical Report (5 pages)", "Stack, design, testing, GenAI"),
        ("127 Automated Tests", "All passing (~4 seconds)"),
        ("Version Control (Git)", "Full commit history"),
        ("GenAI Declaration", "In report & dev log"),
        ("Presentation Slides", "This presentation"),
        ("AI Development Log", "19 entries"),
        ("README.md", "Setup, API reference, deployment"),
    ]

    rows_count = len(deliverables) + 1
    table_shape = slide.shapes.add_table(rows_count, 2,
                                          Inches(2), Inches(1.4), Inches(9), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(5.0)
    table.columns[1].width = Inches(4.0)

    for j, hdr in enumerate(["Deliverable", "Status"]):
        cell = table.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    for i, (name, status) in enumerate(deliverables):
        for j, val in enumerate([f"✅  {name}", status]):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ALT if i % 2 == 0 else BG_DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = LIGHT_GRAY

    add_title_text(slide, "Thank you! Questions?", Inches(1), Inches(6.5), Inches(11), Inches(0.6),
                   font_size=24, bold=True, color=RGBColor(100, 140, 255), alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────────
    output_path = os.path.join(OUTPUT_DIR, "SP500_Tracker_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation PPTX saved to: {output_path}")
    return output_path


# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    build_technical_report()
    build_presentation()
    print("\n✅ All documents generated successfully!")
